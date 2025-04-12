from smolagents import CodeAgent, tool, LiteLLMModel
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
import os
import markdown
import bs4

load_dotenv()

# Global variable to store the presentation
pres = None
DEFAULT_IMAGES_FOLDER = "photos"

@tool
def parse_markdown_script(file_path: str) -> dict:
    """
    Parses a markdown script and returns a structured dictionary with sections, headings, and content.
    
    Args:
        file_path: Path to the markdown file
        
    Returns:
        A dictionary with structured content from the markdown file
    """
    try:
        with open(file_path, 'r') as file:
            content = file.read()
        
        # Convert markdown to HTML for easier parsing
        html = markdown.markdown(content)
        soup = bs4.BeautifulSoup(html, 'html.parser')
        
        # Extract title
        title = soup.find('h1').text if soup.find('h1') else "Presentation"
        
        # Extract sections based on h2 headings
        sections = []
        current_h2 = None
        current_content = []
        
        for element in soup.find_all(['h2', 'p', 'ul', 'ol']):
            if element.name == 'h2':
                # Save previous section
                if current_h2 is not None:
                    sections.append({
                        'title': current_h2,
                        'content': current_content
                    })
                # Start new section
                current_h2 = element.text
                current_content = []
            else:
                # Add content to current section
                if element.name == 'p':
                    # Remove italic markers
                    text = element.text.replace('*', '')
                    current_content.append({'type': 'paragraph', 'text': text})
                elif element.name in ('ul', 'ol'):
                    list_items = []
                    for li in element.find_all('li'):
                        # Remove italic markers
                        text = li.text.replace('*', '')
                        list_items.append(text)
                    current_content.append({'type': 'list', 'items': list_items})
        
        # Add the last section
        if current_h2 is not None:
            sections.append({
                'title': current_h2,
                'content': current_content
            })
        
        result = {
            'title': title,
            'sections': sections
        }
        
        print(f"Successfully parsed markdown file: {file_path}")
        return result
    except Exception as e:
        error_msg = f"Error parsing markdown file: {str(e)}"
        print(error_msg)
        return {"error": error_msg}

@tool
def list_available_images(folder_path: str = DEFAULT_IMAGES_FOLDER) -> list:
    """
    Lists all image files available in the specified folder.
    
    Args:
        folder_path: Path to the folder containing images
        
    Returns:
        A list of image filenames
    """
    try:
        if not os.path.exists(folder_path):
            return []
        
        image_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp')
        images = [
            f for f in os.listdir(folder_path) 
            if os.path.isfile(os.path.join(folder_path, f)) and 
            f.lower().endswith(image_extensions)
        ]
        
        print(f"Found {len(images)} images in {folder_path}")
        return images
    except Exception as e:
        error_msg = f"Error listing images: {str(e)}"
        print(error_msg)
        return []

@tool
def find_relevant_image(keyword: str, available_images: list) -> str:
    """
    Finds the most relevant image for a given keyword from available images.
    
    Args:
        keyword: Keyword to match against image filenames
        available_images: List of available image filenames
        
    Returns:
        The most relevant image filename or empty string if none found
    """
    if not available_images:
        return ""
    
    # Normalize keyword for matching
    keyword = keyword.lower().replace(' ', '_')
    
    # Check for direct matches first
    for image in available_images:
        image_name = os.path.splitext(image)[0].lower()
        if keyword in image_name:
            return image
    
    # If no direct match, try partial matches
    for image in available_images:
        image_name = os.path.splitext(image)[0].lower()
        words = keyword.split('_')
        for word in words:
            if len(word) > 3 and word in image_name:  # Only consider words longer than 3 chars
                return image
    
    # If no match found, return the first image
    return available_images[0]

@tool
def create_presentation() -> str:
    """
    Creates a new blank PowerPoint presentation.
    
    Returns:
        A success message.
    """
    global pres
    pres = Presentation()
    print("Created a new PowerPoint presentation")
    return "Created a new blank PowerPoint presentation."

@tool
def set_presentation_theme(theme_name: str = "professional") -> str:
    """
    Sets the theme for the presentation.
    
    Args:
        theme_name: The theme to apply (professional, creative, minimal)
        
    Returns:
        A success message
    """
    global pres
    if not pres:
        create_presentation()
    
    # Define theme colors
    themes = {
        "professional": {
            "title_color": RGBColor(31, 73, 125),  # Dark blue
            "subtitle_color": RGBColor(68, 114, 196),  # Medium blue
            "text_color": RGBColor(0, 0, 0),  # Black
            "background_color": RGBColor(255, 255, 255),  # White
        },
        "creative": {
            "title_color": RGBColor(192, 0, 0),  # Dark red
            "subtitle_color": RGBColor(255, 102, 0),  # Orange
            "text_color": RGBColor(51, 51, 51),  # Dark gray
            "background_color": RGBColor(248, 248, 248),  # Light gray
        },
        "minimal": {
            "title_color": RGBColor(0, 0, 0),  # Black
            "subtitle_color": RGBColor(100, 100, 100),  # Medium gray
            "text_color": RGBColor(64, 64, 64),  # Dark gray
            "background_color": RGBColor(255, 255, 255),  # White
        }
    }
    
    # If theme doesn't exist, default to professional
    if theme_name not in themes:
        theme_name = "professional"
    
    # Store theme in global variable
    globals()["theme"] = themes[theme_name]
    
    return f"Set presentation theme to {theme_name}"

@tool
def add_title_slide(title: str, subtitle: str = None) -> str:
    """
    Adds a title slide to the presentation.
    
    Args:
        title: The main title text for the slide.
        subtitle: Optional subtitle text for the slide.
        
    Returns:
        A success message.
    """
    global pres
    if not pres:
        create_presentation()
        set_presentation_theme()
        
    slide_layout = pres.slide_layouts[0]  # Title slide layout
    slide = pres.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Apply theme formatting to title
    title_frame = title_shape.text_frame
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(40)
            run.font.bold = True
            if "theme" in globals():
                run.font.color.rgb = globals()["theme"]["title_color"]
    
    # Set subtitle if provided
    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        
        # Apply theme formatting to subtitle
        subtitle_frame = subtitle_shape.text_frame
        for paragraph in subtitle_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(24)
                if "theme" in globals():
                    run.font.color.rgb = globals()["theme"]["subtitle_color"]
        
    return f"Added title slide with title: '{title}'" + (f" and subtitle: '{subtitle}'" if subtitle else "")

@tool
def add_section_slide(title: str) -> str:
    """
    Adds a section slide to the presentation.
    
    Args:
        title: The section title text for the slide.
        
    Returns:
        A success message.
    """
    global pres
    if not pres:
        create_presentation()
        set_presentation_theme()
        
    slide_layout = pres.slide_layouts[2]  # Section header layout
    slide = pres.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Apply theme formatting to title
    title_frame = title_shape.text_frame
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(36)
            run.font.bold = True
            if "theme" in globals():
                run.font.color.rgb = globals()["theme"]["title_color"]
        
    return f"Added section slide with title: '{title}'"

@tool
def add_content_slide(title: str, content_items: list, include_image: bool = False, image_path: str = "") -> str:
    """
    Adds a content slide with title and bullet points or paragraphs.
    
    Args:
        title: The title text for the slide.
        content_items: List of content items (text or bullet points)
        include_image: Whether to include an image
        image_path: Path to the image file (if include_image is True)
        
    Returns:
        A success message.
    """
    global pres
    if not pres:
        create_presentation()
        set_presentation_theme()
        
    # Choose appropriate layout
    if include_image and image_path:
        slide_layout = pres.slide_layouts[8]  # Title, Content and Picture layout
    else:
        slide_layout = pres.slide_layouts[1]  # Title and Content layout
        
    slide = pres.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Apply theme formatting to title
    title_frame = title_shape.text_frame
    for paragraph in title_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            if "theme" in globals():
                run.font.color.rgb = globals()["theme"]["title_color"]
    
    # Add content
    content_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Content placeholder
            content_placeholder = shape
            break
    
    if not content_placeholder:
        return f"Error: Could not find content placeholder on slide"
    
    text_frame = content_placeholder.text_frame
    text_frame.clear()  # Clear existing content
    
    # Add content items
    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
            
        p.text = item
        p.level = 0  # First level bullet
        
        # Apply theme formatting
        for run in p.runs:
            run.font.size = Pt(24)
            if "theme" in globals():
                run.font.color.rgb = globals()["theme"]["text_color"]
    
    # Add image if requested
    if include_image and image_path:
        if os.path.exists(image_path):
            # Find picture placeholder or add picture directly
            picture_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 2:  # Picture placeholder
                    picture_placeholder = shape
                    break
            
            if picture_placeholder:
                # If we have a picture placeholder, use it
                picture_placeholder.insert_picture(image_path)
            else:
                # Otherwise add picture directly to slide
                slide.shapes.add_picture(image_path, Inches(6), Inches(2), width=Inches(4))
        else:
            print(f"Warning: Image path '{image_path}' does not exist.")
    
    return f"Added content slide with title: '{title}' and {len(content_items)} content items"

@tool
def convert_script_to_presentation(script_data: dict, available_images: list) -> str:
    """
    Converts parsed script data into a presentation.
    
    Args:
        script_data: Structured script data from parse_markdown_script
        available_images: List of available image filenames
        
    Returns:
        A success message with slides created
    """
    global pres
    if not pres:
        create_presentation()
        set_presentation_theme("professional")
    
    slides_created = 0
    
    # Add title slide
    title = script_data.get('title', 'Presentation')
    add_title_slide(title, "Created with SmolaGents")
    slides_created += 1
    
    # Process each section
    for section in script_data.get('sections', []):
        section_title = section.get('title', '')
        
        # Add section slide
        add_section_slide(section_title)
        slides_created += 1
        
        # Find relevant image for this section
        image_filename = ""
        if available_images:
            image_filename = find_relevant_image(section_title, available_images)
            if image_filename:
                image_filename = os.path.join(DEFAULT_IMAGES_FOLDER, image_filename)
        
        # Process content for this section
        for content_block in section.get('content', []):
            if content_block.get('type') == 'list':
                # For bullet point lists, create a content slide
                bullet_points = content_block.get('items', [])
                if bullet_points:
                    add_content_slide(
                        section_title, 
                        bullet_points,
                        include_image=bool(image_filename),
                        image_path=image_filename
                    )
                    slides_created += 1
                    # Only use image once per section
                    image_filename = ""
            
            elif content_block.get('type') == 'paragraph':
                # For paragraphs, create content slides with the paragraph text split into bullet points
                paragraph_text = content_block.get('text', '')
                if paragraph_text:
                    # Split long paragraphs into bullet points
                    sentences = re.split(r'(?<=[.!?])\s+', paragraph_text)
                    bullet_points = []
                    current_point = ""
                    
                    for sentence in sentences:
                        if sentence:
                            if len(current_point + sentence) > 100:  # Limit bullet point length
                                if current_point:
                                    bullet_points.append(current_point)
                                current_point = sentence
                            else:
                                current_point += " " + sentence if current_point else sentence
                    
                    if current_point:
                        bullet_points.append(current_point)
                    
                    if bullet_points:
                        add_content_slide(
                            section_title, 
                            bullet_points,
                            include_image=bool(image_filename),
                            image_path=image_filename
                        )
                        slides_created += 1
                        # Only use image once per section
                        image_filename = ""
    
    return f"Successfully created presentation with {slides_created} slides based on the script."

@tool
def save_presentation(filename: str) -> str:
    """
    Saves the presentation to a file.
    
    Args:
        filename: The filename to save the presentation to (will add .pptx extension if missing).
        
    Returns:
        A success or error message.
    """
    global pres
    if not pres:
        return "Error: No presentation has been created yet. Use create_presentation first."
        
    if not filename.endswith('.pptx'):
        filename += '.pptx'
        
    pres.save(filename)
    return f"Presentation saved to '{filename}'"

# Initialize model with Gemini 2.0 Flash
model = LiteLLMModel(
    model_id="gemini/gemini-2.0-flash",
    temperature=0.2,
    api_key=os.environ["GEMINI_API_KEY"],
)

# Create the agent
script_to_presentation_agent = CodeAgent(
    tools=[
        parse_markdown_script,
        list_available_images,
        find_relevant_image,
        create_presentation,
        set_presentation_theme,
        add_title_slide,
        add_section_slide,
        add_content_slide,
        convert_script_to_presentation,
        save_presentation
    ],
    model=model,
    add_base_tools=True,
    verbosity_level=2
)

if __name__ == "__main__":
    # Example usage
    result = script_to_presentation_agent.run(
        """
        Convert the YouTube script in 'yt_script.md' into a polished PowerPoint presentation.
        The presentation should:
        1. Have a professional theme
        2. Include title slides for each section
        3. Break down paragraphs into bullet points
        4. Include images where available and relevant
        5. Save the final presentation as 'ai_agent_presentation.pptx'
        
        Make the presentation visually appealing and professional, suitable for a tech talk or conference presentation.
        """
    )
    print(result)